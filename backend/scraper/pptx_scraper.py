"""
PPTX Scraper for extracting and ingesting PowerPoint documents
"""
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from scraper.document_processor import DocumentProcessor
from services.embedding_service import embedding_service

logger = logging.getLogger(__name__)


class PPTXScraper:
    """Scraper for PowerPoint (PPTX) documents"""
    
    def __init__(self):
        """Initialize PPTX scraper"""
        self.document_processor = DocumentProcessor()
        
        if not PPTX_AVAILABLE:
            logger.warning(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
    
    def extract_text_from_slide(self, slide) -> str:
        """
        Extract text from a single slide
        
        Args:
            slide: Slide object
            
        Returns:
            Extracted text
        """
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_content.append(shape.text.strip())
            
            # Extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text:
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
        
        return "\n".join(text_content)
    
    def extract_text(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract text from PPTX file
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Dictionary with full text and slide-by-slide content
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        try:
            prs = Presentation(pptx_path)
            
            slides_content = []
            full_text = []
            
            logger.info(f"Extracting text from {len(prs.slides)} slides")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = self.extract_text_from_slide(slide)
                
                if slide_text:
                    slides_content.append({
                        'slide_number': slide_num,
                        'content': slide_text
                    })
                    full_text.append(f"Slide {slide_num}:\n{slide_text}")
                    logger.debug(f"Extracted {len(slide_text)} characters from slide {slide_num}")
            
            result = {
                'full_text': '\n\n'.join(full_text),
                'slides': slides_content,
                'total_slides': len(prs.slides),
                'slides_with_content': len(slides_content)
            }
            
            logger.info(
                f"Extracted text from {result['slides_with_content']} of "
                f"{result['total_slides']} slides"
            )
            
            return result
            
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            raise
    
    def extract_metadata(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract metadata from PPTX file
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Metadata dictionary
        """
        metadata = {
            'filename': Path(pptx_path).name,
            'filepath': str(pptx_path),
            'file_size': Path(pptx_path).stat().st_size,
            'source_type': 'PowerPoint Presentation',
            'ingestion_date': datetime.utcnow().isoformat()
        }
        
        if PPTX_AVAILABLE:
            try:
                prs = Presentation(pptx_path)
                
                metadata['total_slides'] = len(prs.slides)
                
                # Extract core properties
                core_props = prs.core_properties
                metadata['title'] = core_props.title or metadata['filename']
                metadata['author'] = core_props.author or 'Unknown'
                metadata['subject'] = core_props.subject or ''
                metadata['keywords'] = core_props.keywords or ''
                metadata['comments'] = core_props.comments or ''
                metadata['category'] = core_props.category or ''
                
                if core_props.created:
                    metadata['created_date'] = core_props.created.isoformat()
                if core_props.modified:
                    metadata['modified_date'] = core_props.modified.isoformat()
                
            except Exception as e:
                logger.warning(f"Could not extract PPTX metadata: {str(e)}")
                metadata['title'] = metadata['filename']
        else:
            metadata['title'] = metadata['filename']
        
        return metadata
    
    async def process_pptx(
        self,
        pptx_path: str,
        namespace: str = "",
        custom_metadata: Optional[Dict[str, Any]] = None,
        chunk_by_slide: bool = False
    ) -> Dict[str, Any]:
        """
        Process a PPTX file and ingest into Pinecone
        
        Args:
            pptx_path: Path to PPTX file
            namespace: Pinecone namespace
            custom_metadata: Additional metadata
            chunk_by_slide: If True, create one chunk per slide; if False, use standard chunking
            
        Returns:
            Processing results
        """
        try:
            logger.info(f"Processing PPTX: {pptx_path}")
            
            # Check if file exists
            if not Path(pptx_path).exists():
                raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
            
            # Extract text
            extraction_result = self.extract_text(pptx_path)
            full_text = extraction_result['full_text']
            
            if not full_text or len(full_text.strip()) < 100:
                logger.warning(f"Extracted text too short ({len(full_text)} chars), skipping")
                return {
                    'success': False,
                    'error': 'Extracted text too short or empty',
                    'file': pptx_path
                }
            
            # Extract metadata
            metadata = self.extract_metadata(pptx_path)
            
            # Add custom metadata
            if custom_metadata:
                metadata.update(custom_metadata)
            
            # Create chunks
            if chunk_by_slide:
                # One chunk per slide
                chunks = []
                for slide_data in extraction_result['slides']:
                    slide_metadata = metadata.copy()
                    slide_metadata.update({
                        'slide_number': slide_data['slide_number'],
                        'total_slides': extraction_result['total_slides']
                    })
                    
                    chunks.append({
                        'content': slide_data['content'],
                        'metadata': slide_metadata
                    })
            else:
                # Standard chunking
                document = {
                    'content': full_text,
                    'metadata': metadata
                }
                chunks = self.document_processor.create_chunks_with_metadata(document)
            
            if not chunks:
                logger.warning(f"No chunks created from PPTX: {pptx_path}")
                return {
                    'success': False,
                    'error': 'No chunks created',
                    'file': pptx_path
                }
            
            logger.info(f"Created {len(chunks)} chunks from PPTX")
            
            # Ingest into Pinecone
            results = await embedding_service.ingest_documents(
                chunks,
                namespace=namespace
            )
            
            return {
                'success': True,
                'file': pptx_path,
                'chunks_created': len(chunks),
                'vectors_stored': results.get('stored', 0),
                'slides_processed': extraction_result['slides_with_content'],
                'total_slides': extraction_result['total_slides'],
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing PPTX {pptx_path}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'file': pptx_path
            }
    
    async def process_pptx_directory(
        self,
        directory_path: str,
        namespace: str = "",
        recursive: bool = True,
        custom_metadata: Optional[Dict[str, Any]] = None,
        chunk_by_slide: bool = False
    ) -> Dict[str, Any]:
        """
        Process all PPTX files in a directory
        
        Args:
            directory_path: Path to directory
            namespace: Pinecone namespace
            recursive: Search subdirectories
            custom_metadata: Additional metadata for all files
            chunk_by_slide: Chunk by slide or use standard chunking
            
        Returns:
            Processing results summary
        """
        try:
            directory = Path(directory_path)
            
            if not directory.exists():
                raise FileNotFoundError(f"Directory not found: {directory_path}")
            
            # Find all PPTX files
            if recursive:
                pptx_files = list(directory.rglob('*.pptx')) + list(directory.rglob('*.ppt'))
            else:
                pptx_files = list(directory.glob('*.pptx')) + list(directory.glob('*.ppt'))
            
            logger.info(f"Found {len(pptx_files)} PowerPoint files in {directory_path}")
            
            if not pptx_files:
                return {
                    'success': True,
                    'message': 'No PowerPoint files found',
                    'total_files': 0,
                    'processed': 0,
                    'failed': 0
                }
            
            results = {
                'total_files': len(pptx_files),
                'processed': 0,
                'failed': 0,
                'total_chunks': 0,
                'total_vectors': 0,
                'total_slides': 0,
                'files': []
            }
            
            # Process each PPTX
            for pptx_file in pptx_files:
                logger.info(f"Processing {pptx_file.name}...")
                
                result = await self.process_pptx(
                    str(pptx_file),
                    namespace=namespace,
                    custom_metadata=custom_metadata,
                    chunk_by_slide=chunk_by_slide
                )
                
                if result['success']:
                    results['processed'] += 1
                    results['total_chunks'] += result.get('chunks_created', 0)
                    results['total_vectors'] += result.get('vectors_stored', 0)
                    results['total_slides'] += result.get('total_slides', 0)
                else:
                    results['failed'] += 1
                
                results['files'].append(result)
            
            logger.info(
                f"Completed processing: {results['processed']} successful, "
                f"{results['failed']} failed, {results['total_vectors']} vectors stored"
            )
            
            return {
                'success': True,
                **results
            }
            
        except Exception as e:
            logger.error(f"Error processing directory {directory_path}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'directory': directory_path
            }


# Global instance
pptx_scraper = PPTXScraper()

# Made with Bob